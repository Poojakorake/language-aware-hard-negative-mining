"""
Generates presentation slides for:
"Language-Aware Hard Negative Mining for Improved Vision-Language Representation Learning"
CS 6384 Computer Vision Final Project — UTD
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Colour palette ────────────────────────────────────────────────────────────
UTD_ORANGE  = RGBColor(0xC7, 0x5B, 0x12)   # UTD brand orange
UTD_GREEN   = RGBColor(0x15, 0x4F, 0x37)   # UTD brand green  (unused but available)
DARK_GRAY   = RGBColor(0x2E, 0x2E, 0x2E)
MID_GRAY    = RGBColor(0x55, 0x55, 0x55)
LIGHT_GRAY  = RGBColor(0xF2, 0xF2, 0xF2)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT_BLUE = RGBColor(0x1F, 0x6F, 0xB8)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

blank_layout = prs.slide_layouts[6]   # completely blank


# ── Helper functions ──────────────────────────────────────────────────────────

def add_rect(slide, l, t, w, h, fill_color, border_color=None, border_pt=0):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color and border_pt:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(border_pt)
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, l, t, w, h,
             font_size=18, bold=False, italic=False,
             color=DARK_GRAY, align=PP_ALIGN.LEFT,
             wrap=True, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txBox.word_wrap = wrap
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font_name
    return txBox


def add_bullet_box(slide, items, l, t, w, h,
                   font_size=17, color=DARK_GRAY,
                   bullet="▸  ", line_spacing=1.15):
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txBox.word_wrap = True
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(4)
        run = p.add_run()
        run.text = bullet + item
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        run.font.name = "Calibri"


def header_bar(slide, title, subtitle=None):
    """Orange bar at top with slide title."""
    add_rect(slide, 0, 0, 13.33, 1.1, UTD_ORANGE)
    add_text(slide, title, 0.35, 0.1, 12.5, 0.75,
             font_size=30, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        add_text(slide, subtitle, 0.35, 0.75, 12.5, 0.35,
                 font_size=15, italic=True, color=WHITE, align=PP_ALIGN.LEFT)
    # thin bottom accent line
    add_rect(slide, 0, 1.1, 13.33, 0.04, DARK_GRAY)


def footer(slide, slide_num, total=9):
    add_rect(slide, 0, 7.2, 13.33, 0.3, DARK_GRAY)
    add_text(slide, "CS 6384 Computer Vision  |  UTD  |  Spring 2025",
             0.2, 7.2, 10, 0.3, font_size=10, color=WHITE, align=PP_ALIGN.LEFT)
    add_text(slide, f"{slide_num} / {total}",
             12.5, 7.2, 0.8, 0.3, font_size=10, color=WHITE, align=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════════════════════════
# Slide 1 — Title
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)

add_rect(slide, 0, 0, 13.33, 7.5, DARK_GRAY)             # dark background
add_rect(slide, 0, 0, 13.33, 0.08, UTD_ORANGE)           # top accent
add_rect(slide, 0, 7.42, 13.33, 0.08, UTD_ORANGE)        # bottom accent

add_text(slide,
    "Language-Aware Hard Negative Mining\nfor Improved Vision-Language\nRepresentation Learning",
    0.8, 1.2, 11.8, 3.0,
    font_size=38, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_rect(slide, 3.5, 4.1, 6.33, 0.05, UTD_ORANGE)        # divider line

add_text(slide,
    "Prafull Kumar Prajapati  ·  Pooja Ramesh Korake  ·  Faizul Anis",
    0.5, 4.3, 12.3, 0.5,
    font_size=18, color=RGBColor(0xCC, 0xCC, 0xCC), align=PP_ALIGN.CENTER)

add_text(slide,
    "The University of Texas at Dallas",
    0.5, 4.9, 12.3, 0.5,
    font_size=16, italic=True, color=UTD_ORANGE, align=PP_ALIGN.CENTER)

add_text(slide,
    "CS 6384 Computer Vision  |  Spring 2025",
    0.5, 5.5, 12.3, 0.5,
    font_size=15, color=RGBColor(0xAA, 0xAA, 0xAA), align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# Slide 2 — Motivation
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
header_bar(slide, "Motivation", "Why does hard negative mining matter?")
footer(slide, 2)

add_text(slide, "Problem with standard contrastive learning:",
         0.4, 1.3, 12, 0.5, font_size=20, bold=True, color=UTD_ORANGE)

add_bullet_box(slide, [
    "CLIP-style training pushes matched (image, text) pairs together and unmatched pairs apart.",
    "Easy negatives (very different images/captions) provide little gradient signal.",
    "Hard negatives — pairs that are similar but not matched — are crucial for learning fine-grained distinctions.",
], 0.5, 1.85, 12.3, 1.8, font_size=18)

add_rect(slide, 0.4, 3.55, 12.5, 0.04, UTD_ORANGE)

add_text(slide, "The LLaVE insight (our starting point):",
         0.4, 3.7, 12, 0.5, font_size=20, bold=True, color=UTD_ORANGE)

add_bullet_box(slide, [
    "Dynamically reweight negatives by their embedding-space similarity (hardness score).",
    "Removes the need for separate offline mining stages.",
    "BUT: hardness computed purely from the model's own embeddings has a flaw...",
], 0.5, 4.25, 12.3, 1.8, font_size=18)

add_rect(slide, 0.5, 6.05, 12.3, 0.7, RGBColor(0xFF, 0xF3, 0xE0))
add_text(slide,
    '⚠  Two captions saying the same thing differently (e.g. "A dog runs" vs "A canine sprints") '
    'get treated as hard negatives — wasting model capacity.',
    0.6, 6.1, 12.1, 0.6, font_size=15, italic=True, color=DARK_GRAY)


# ══════════════════════════════════════════════════════════════════════════════
# Slide 3 — Our Approach
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
header_bar(slide, "Our Approach", "Language-Aware Hardness Score")
footer(slide, 3)

add_text(slide, "Key Idea: inject an external linguistic prior (Sentence-BERT)",
         0.4, 1.25, 12.5, 0.5, font_size=20, bold=True, color=DARK_GRAY)

# Equation box
add_rect(slide, 1.2, 1.85, 10.9, 1.1, RGBColor(0xE8, 0xF4, 0xFF))
add_text(slide,
    "H(i,j)  =  α · sim_embed(i,j)  +  β · sim_ling(i,j)",
    1.3, 1.95, 10.7, 0.6,
    font_size=24, bold=True, color=ACCENT_BLUE, align=PP_ALIGN.CENTER)
add_text(slide, "w(i,j) = exp( H(i,j) )",
    1.3, 2.55, 10.7, 0.4,
    font_size=20, bold=False, color=ACCENT_BLUE, align=PP_ALIGN.CENTER)

# Two columns
add_text(slide, "sim_embed", 0.5, 3.2, 5.5, 0.4, font_size=18, bold=True, color=UTD_ORANGE)
add_bullet_box(slide, [
    "Cosine similarity in the model's embedding space",
    "Scaled by logit temperature τ = 50",
    "Updates each training step as model learns",
], 0.5, 3.65, 5.8, 1.8, font_size=16)

add_text(slide, "sim_ling  (NEW)", 7.0, 3.2, 5.8, 0.4, font_size=18, bold=True, color=UTD_ORANGE)
add_bullet_box(slide, [
    "Cosine similarity from frozen Sentence-BERT",
    "Pre-computed once — zero overhead at train time",
    "Catches semantically equivalent captions",
], 7.0, 3.65, 5.8, 1.8, font_size=16)

add_rect(slide, 6.5, 3.1, 0.04, 2.8, MID_GRAY)   # vertical divider

add_rect(slide, 0.4, 5.6, 12.5, 0.85, LIGHT_GRAY)
add_text(slide,
    "β = 0  →  standard LLaVE baseline     |     α=5, β=5  →  our best (Balanced) configuration\n"
    "Ablation: (α,β) ∈ {(9,0), (9,1), (5,5), (9,5), (1,9)}  — equal weighting wins",
    0.5, 5.65, 12.3, 0.75, font_size=16, color=DARK_GRAY, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# Slide 4 — Method Overview (diagram)
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
header_bar(slide, "Method Overview")
footer(slide, 4)

# Flow diagram using shapes
boxes = [
    (0.4,  3.2, 2.5, 0.85, "Image + Text\nPairs", ACCENT_BLUE),
    (3.3,  3.2, 2.5, 0.85, "CLIP Encoder\n(trainable)", UTD_ORANGE),
    (6.2,  2.55, 2.5, 0.85, "sim_embed\n(embedding space)", ACCENT_BLUE),
    (6.2,  3.9,  2.5, 0.85, "sim_ling\n(Sentence-BERT, frozen)", RGBColor(0x2E, 0x7D, 0x32)),
    (9.3,  3.2, 2.5, 0.85, "Hardness Score\nH = α·E + β·L", UTD_ORANGE),
    (9.3,  4.45, 2.5, 0.85, "Weighted\nInfoNCE Loss", DARK_GRAY),
]
for (l, t, w, h, txt, col) in boxes:
    add_rect(slide, l, t, w, h, col)
    add_text(slide, txt, l+0.05, t+0.05, w-0.1, h-0.1,
             font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Arrows (thin orange rectangles)
add_rect(slide, 2.9,  3.58, 0.4, 0.08, UTD_ORANGE)   # box1 -> box2
add_rect(slide, 5.8,  3.0,  0.4, 0.08, UTD_ORANGE)   # box2 -> sim_embed
add_rect(slide, 5.8,  4.25, 0.4, 0.08, UTD_ORANGE)   # box2 -> sim_ling
add_rect(slide, 8.8,  3.58, 0.5, 0.08, UTD_ORANGE)   # sims -> H
add_rect(slide, 10.55,4.05, 0.08, 0.4, UTD_ORANGE)   # H -> loss (vertical)

# Sentence-BERT pre-compute note
add_rect(slide, 0.4, 5.7, 5.5, 0.75, RGBColor(0xE8, 0xF5, 0xE9))
add_text(slide,
    "Pre-compute step (simularity.py):\nAll captions → Sentence-BERT → .pt embeddings saved to disk",
    0.5, 5.72, 5.3, 0.7, font_size=13, color=RGBColor(0x2E, 0x7D, 0x32))

add_rect(slide, 6.2, 5.7, 5.5, 0.75, RGBColor(0xFF, 0xF3, 0xE0))
add_text(slide,
    "At training time: batch indices → O(1) lookup of pre-computed\nSBERT embeddings. Zero model parameters added.",
    6.3, 5.72, 5.3, 0.7, font_size=13, color=RGBColor(0xBF, 0x36, 0x00))


# ══════════════════════════════════════════════════════════════════════════════
# Slide 5 — Numerical Stability
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
header_bar(slide, "Implementation Details", "Making it work in BFloat16")
footer(slide, 5)

add_text(slide, "Two engineering fixes were needed for stable training:",
         0.4, 1.3, 12.5, 0.45, font_size=19, bold=True, color=DARK_GRAY)

# Fix 1
add_rect(slide, 0.4, 1.9, 0.45, 1.5, UTD_ORANGE)
add_text(slide, "Fix 1", 0.45, 1.95, 0.4, 0.4, font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(slide, "FP32 Upcast", 1.0, 1.9, 5.5, 0.45, font_size=18, bold=True, color=UTD_ORANGE)
add_bullet_box(slide, [
    "BFloat16 rounds small values to 0 — loss collapses to 0/0",
    "Cast sim_embed and sim_ling to float32 before any exponential",
    "Prevents NaN gradients during training",
], 1.0, 2.35, 11.5, 1.1, font_size=16)

# Fix 2
add_rect(slide, 0.4, 3.65, 0.45, 1.5, UTD_ORANGE)
add_text(slide, "Fix 2", 0.45, 3.7, 0.4, 0.4, font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(slide, "Log-Sum-Exp Trick", 1.0, 3.65, 5.5, 0.45, font_size=18, bold=True, color=UTD_ORANGE)
add_bullet_box(slide, [
    "Naively computing exp(sim) creates huge intermediate values",
    "Subtract row maximum before exponentiation — numerically identical but stable",
    "Entire loss stays in log-space; no intermediate 1.0/1.0 rounding",
], 1.0, 4.1, 11.5, 1.2, font_size=16)

add_rect(slide, 0.4, 5.3, 12.5, 1.0, LIGHT_GRAY)
add_text(slide,
    "Pre-computation:  simularity.py encodes all captions with frozen Sentence-BERT once before training.\n"
    "Result: (M × 384) tensor saved to disk.  Training: batch indices → O(1) lookup. No GPU overhead.",
    0.55, 5.35, 12.3, 0.9, font_size=15, color=DARK_GRAY)


# ══════════════════════════════════════════════════════════════════════════════
# Slide 6 — Experimental Setup
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
header_bar(slide, "Experimental Setup")
footer(slide, 6)

# Left column
add_text(slide, "Dataset", 0.4, 1.3, 5.9, 0.45, font_size=19, bold=True, color=UTD_ORANGE)
add_bullet_box(slide, [
    "Flickr30K — 1,000 test images (Karpathy split)",
    "1 caption per image in evaluation",
    "1,000-way retrieval: both I→T and T→I",
], 0.4, 1.75, 6.0, 1.3, font_size=17)

add_text(slide, "Models", 0.4, 3.2, 5.9, 0.45, font_size=19, bold=True, color=UTD_ORANGE)
add_bullet_box(slide, [
    "Backbone: LLaVE-2B (LlavaQwenForCausalLM + SigLIP-so400m-patch14-384)",
    "Linguistic prior: frozen Sentence-BERT all-MiniLM-L6-v2",
    "AdamW lr=2e-4, batch 4 × accum 8 = eff. 32, 3 epochs, DeepSpeed ZeRO-2",
], 0.4, 3.65, 6.0, 1.3, font_size=17)

add_text(slide, "Hyperparameters Ablated", 0.4, 5.1, 5.9, 0.45, font_size=19, bold=True, color=UTD_ORANGE)
add_bullet_box(slide, [
    "(α=9, β=0) baseline  |  (α=9, β=1) language-heavy",
    "(α=5, β=5) balanced ★  |  (α=9, β=5) high αβ",
    "(α=1, β=9) embed-heavy  — 5 configs total",
], 0.4, 5.55, 6.0, 1.3, font_size=17)

# Right column — metrics box
add_rect(slide, 6.8, 1.3, 6.1, 5.5, LIGHT_GRAY)
add_text(slide, "Evaluation Metrics", 7.0, 1.4, 5.7, 0.45, font_size=19, bold=True, color=DARK_GRAY)
add_bullet_box(slide, [
    "Recall@1  — correct item in top 1",
    "Recall@5  — correct item in top 5",
    "Recall@10 — correct item in top 10",
    "Median Rank (MedR) — lower is better",
    "",
    "Reported for both directions:",
    "  Image → Text  (I→T)",
    "  Text → Image  (T→I)",
], 7.0, 1.9, 5.7, 4.5, font_size=17)


# ══════════════════════════════════════════════════════════════════════════════
# Slide 7 — Results
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
header_bar(slide, "Results on Flickr30K", "T→I retrieval — 1K images, 5K captions (evaluate_retrieval.py)")
footer(slide, 7)

# Table header
add_rect(slide, 0.3, 1.3, 12.7, 0.5, DARK_GRAY)
for txt, l in [("Configuration (α, β)", 0.4), ("R@1", 5.2), ("R@5", 6.8), ("R@10", 8.4), ("MedR", 10.5)]:
    add_text(slide, txt, l, 1.33, 2.2, 0.45, font_size=14, bold=True,
             color=WHITE, align=PP_ALIGN.CENTER)

# Table rows
rows = [
    ("Baseline (9, 0)",       "6.70",  "17.34", "25.16", "47",  WHITE,                          DARK_GRAY),
    ("Language-Heavy (9, 1)", "8.64",  "22.00", "30.50", "33",  LIGHT_GRAY,                     DARK_GRAY),
    ("Balanced (5, 5) ★ BEST","8.70",  "22.26", "30.64", "33",  RGBColor(0xFF, 0xF0, 0xD8),     UTD_ORANGE),
    ("High αβ (9, 5)",        "8.04",  "21.00", "28.80", "37",  LIGHT_GRAY,                     DARK_GRAY),
    ("Embed-Heavy (1, 9)",    "7.86",  "20.84", "28.78", "36",  WHITE,                          DARK_GRAY),
]
for i, (name, r1, r5, r10, medr, bg, fc) in enumerate(rows):
    t = 1.8 + i * 0.6
    add_rect(slide, 0.3, t, 12.7, 0.58, bg)
    add_text(slide, name, 0.4, t+0.08, 4.7, 0.45, font_size=14,
             bold=(i==2), color=fc, align=PP_ALIGN.LEFT)
    for val, l in [(r1, 5.2), (r5, 6.8), (r10, 8.4), (medr, 10.5)]:
        add_text(slide, val, l, t+0.08, 1.5, 0.45, font_size=14,
                 bold=(i==2), color=fc, align=PP_ALIGN.CENTER)

# Gain row
add_rect(slide, 0.3, 4.8, 12.7, 0.45, RGBColor(0xE8, 0xF4, 0xFF))
for txt, l in [("vs Baseline ↑", 0.4), ("+2.0 pp", 5.2), ("+4.92 pp", 6.8), ("+5.48 pp (+21.8% rel)", 7.8), ("47→33", 10.5)]:
    add_text(slide, txt, l, 4.83, 2.5, 0.42, font_size=13, italic=True,
             color=ACCENT_BLUE, align=PP_ALIGN.CENTER)

# Key takeaways
add_text(slide, "Key Takeaways:", 0.4, 5.4, 12.5, 0.4, font_size=17, bold=True, color=DARK_GRAY)
add_bullet_box(slide, [
    "Balanced (α=5, β=5) is best — equal weighting of embedding and linguistic signals.",
    "R@10 improves 25.16% → 30.64% (+21.8% relative); MedR drops from 47 to 33.",
    "High αβ (9,5) underperforms Balanced (5,5) despite same β — gradient conflict.",
], 0.5, 5.85, 12.3, 1.4, font_size=16)


# ══════════════════════════════════════════════════════════════════════════════
# Slide 8 — Ablation
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
header_bar(slide, "Ablation Study", "Effect of (α, β) balance — T→I R@10 on Flickr30K")
footer(slide, 8)

add_text(slide, "T→I R@10 (%) across (α, β) configurations", 2.5, 1.3, 8.0, 0.45,
         font_size=18, bold=True, color=DARK_GRAY, align=PP_ALIGN.CENTER)

bar_data = [
    ("Baseline\n(9,0)",      25.16, DARK_GRAY),
    ("Lang-Heavy\n(9,1)",    30.50, ACCENT_BLUE),
    ("Balanced\n(5,5) ★",   30.64, UTD_ORANGE),
    ("High αβ\n(9,5)",       28.80, MID_GRAY),
    ("Embed-Heavy\n(1,9)",   28.78, MID_GRAY),
]
chart_bottom = 6.5
chart_height = 4.0
max_val = 34.0
min_val = 22.0
span = max_val - min_val
bar_w = 1.5
gap   = 0.35
start_l = 0.55

for i, (label, val, color) in enumerate(bar_data):
    bar_h = (val - min_val) / span * chart_height
    l = start_l + i * (bar_w + gap)
    t = chart_bottom - bar_h
    add_rect(slide, l, t, bar_w, bar_h, color)
    add_text(slide, f"{val}%", l, t - 0.45, bar_w, 0.4,
             font_size=15, bold=(i==2), color=color, align=PP_ALIGN.CENTER)
    add_text(slide, label, l, chart_bottom + 0.05, bar_w, 0.6,
             font_size=13, color=DARK_GRAY, align=PP_ALIGN.CENTER)

# Y-axis grid lines
for v in [24, 26, 28, 30, 32, 34]:
    bar_h = (v - min_val) / span * chart_height
    t = chart_bottom - bar_h
    add_text(slide, f"{v}", 9.7, t - 0.15, 0.7, 0.35,
             font_size=11, color=MID_GRAY, align=PP_ALIGN.RIGHT)
    add_rect(slide, 0.5, t, 9.3, 0.01, RGBColor(0xCC, 0xCC, 0xCC))

# Right side: gradient conflict note
add_rect(slide, 10.2, 1.3, 3.0, 5.6, LIGHT_GRAY)
add_text(slide, "Observations:", 10.3, 1.4, 2.8, 0.4, font_size=16, bold=True, color=DARK_GRAY)
add_bullet_box(slide, [
    "Adding β=1 to baseline already closes most of the gap.",
    "Balanced (5,5) edges out (9,1) — equal weighting is optimal.",
    "High αβ (9,5) drops vs. Balanced: gradient conflict.",
    "Embed-Heavy (1,9): over-weighting text hurts.",
    "",
    "→ Equal α=β=5 provides the most stable gradient signal.",
], 10.3, 1.85, 2.8, 4.5, font_size=13)


# ══════════════════════════════════════════════════════════════════════════════
# Slide 9 — Conclusion
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
header_bar(slide, "Conclusion & Future Work")
footer(slide, 9)

# Two column layout
add_text(slide, "What we did:", 0.4, 1.3, 6.0, 0.45, font_size=20, bold=True, color=UTD_ORANGE)
add_bullet_box(slide, [
    "Extended LLaVE's hardness-weighted contrastive loss with a frozen Sentence-BERT linguistic prior.",
    "Combined embedding similarity and text semantic similarity into one unified hardness score.",
    "Added FP32 upcast + log-sum-exp for BFloat16 training stability.",
    "Best model (α=5, β=5): +21.8% relative R@10 improvement; MedR 47→33.",
], 0.4, 1.8, 6.0, 3.2, font_size=17)

add_rect(slide, 6.7, 1.2, 0.04, 4.5, MID_GRAY)

add_text(slide, "Future directions:", 7.0, 1.3, 6.0, 0.45, font_size=20, bold=True, color=UTD_ORANGE)
add_bullet_box(slide, [
    "Scale to full Flickr30K + MS-COCO training sets (GPU cluster).",
    "Explore multimodal linguistic prior (image + caption jointly).",
    "Learnable α and β rather than fixed hyperparameters.",
    "Apply to cross-lingual retrieval where semantic overlap is harder to detect.",
], 7.0, 1.8, 6.0, 3.2, font_size=17)

# Summary box
add_rect(slide, 0.4, 5.5, 12.5, 1.2, DARK_GRAY)
add_text(slide,
    "Core message:  Equal weighting of embedding + linguistic similarity (α=5, β=5) prevents\n"
    "semantically equivalent captions from being treated as hard negatives —\n"
    "21.8% relative R@10 improvement and MedR 47→33, with no extra learnable parameters.",
    0.6, 5.55, 12.1, 1.1, font_size=16, color=WHITE, align=PP_ALIGN.CENTER)


# ── Save ─────────────────────────────────────────────────────────────────────
out_path = "/Users/poojakorake/Documents/Computer Vision/Final Report/presentation.pptx"
prs.save(out_path)
print(f"Saved: {out_path}")
print(f"Slides: 9")
print("Open in PowerPoint or Google Slides to present.")
